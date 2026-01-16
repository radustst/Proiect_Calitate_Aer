"""
Convert PROJECT_PRESENTATION.md to PowerPoint format
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

def parse_markdown_slides(md_content):
    """Parse markdown content into slides"""
    slides_text = md_content.split('\n---\n')
    slides = []
    
    for slide_text in slides_text:
        lines = slide_text.strip().split('\n')
        slides.append(lines)
    
    return slides

def add_title_slide(prs, lines):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    # Extract title and subtitle
    title_text = ""
    subtitle_text = []
    
    for line in lines:
        line = line.strip()
        if line.startswith('# '):
            title_text = line[2:]
        elif line.startswith('## '):
            subtitle_text.append(line[3:])
        elif line.startswith('**') and line.endswith('**'):
            subtitle_text.append(line[2:-2])
        elif line and not line.startswith('#'):
            subtitle_text.append(line)
    
    title.text = title_text
    subtitle.text = '\n'.join(subtitle_text)
    
    return slide

def add_content_slide(prs, lines):
    """Add a content slide with title and bullet points"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    # Extract title and content
    title_text = ""
    content_lines = []
    
    for line in lines:
        line_stripped = line.strip()
        if line_stripped.startswith('# '):
            title_text = line_stripped[2:]
        elif line_stripped.startswith('## '):
            if not title_text:
                title_text = line_stripped[3:]
            else:
                content_lines.append(('header', line_stripped[3:]))
        elif line_stripped:
            content_lines.append(('text', line_stripped))
    
    title.text = title_text
    
    # Add content
    tf = content.text_frame
    tf.clear()
    
    for i, (content_type, text) in enumerate(content_lines):
        # Remove markdown formatting
        text = text.replace('**', '')
        text = re.sub(r'`([^`]+)`', r'\1', text)
        
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = text
        
        # Set formatting
        if content_type == 'header':
            p.font.bold = True
            p.font.size = Pt(20)
            p.level = 0
        elif text.startswith('- ') or text.startswith('✅ ') or text.startswith('🔹 '):
            # Bullet point
            p.text = text[2:].strip() if text.startswith('- ') else text
            p.level = 1
            p.font.size = Pt(16)
        elif text.startswith('  - '):
            # Sub-bullet
            p.text = text[4:]
            p.level = 2
            p.font.size = Pt(14)
        else:
            p.level = 0
            p.font.size = Pt(18)
    
    return slide

def add_table_slide(prs, lines):
    """Add a slide with a table"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
    )
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    
    # Extract title and table content
    title_text = ""
    table_rows = []
    
    for line in lines:
        line_stripped = line.strip()
        if line_stripped.startswith('# '):
            title_text = line_stripped[2:]
        elif line_stripped.startswith('## '):
            if not title_text:
                title_text = line_stripped[3:]
        elif '|' in line_stripped and not line_stripped.startswith('|---'):
            # Table row
            cells = [cell.strip() for cell in line_stripped.split('|')]
            cells = [c for c in cells if c]  # Remove empty cells
            if cells:
                table_rows.append(cells)
    
    title_para.text = title_text
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    
    # Add table if found
    if table_rows and len(table_rows) > 1:
        rows = len(table_rows)
        cols = max(len(row) for row in table_rows)
        
        left = Inches(1.5)
        top = Inches(2)
        width = Inches(6)
        height = Inches(3)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Fill table
        for i, row_data in enumerate(table_rows):
            for j, cell_text in enumerate(row_data):
                cell = table.rows[i].cells[j]
                cell.text = cell_text.replace('**', '')
                
                # Header row formatting
                if i == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = RGBColor(255, 255, 255)
                        paragraph.font.size = Pt(14)
                else:
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(12)
    
    return slide

def convert_md_to_pptx(md_file, output_file):
    """Convert markdown file to PowerPoint presentation"""
    # Read markdown file
    with open(md_file, 'r', encoding='utf-8') as f:
        md_content = f.read()
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Parse slides
    slides_data = parse_markdown_slides(md_content)
    
    # Add slides
    for i, lines in enumerate(slides_data):
        # Check if it's a title slide (first slide)
        if i == 0:
            add_title_slide(prs, lines)
        # Check if it contains a table
        elif any('|' in line for line in lines):
            add_table_slide(prs, lines)
        else:
            add_content_slide(prs, lines)
    
    # Save presentation
    prs.save(output_file)
    print(f"✅ Presentation saved to: {output_file}")

if __name__ == "__main__":
    md_file = "PROJECT_PRESENTATION.md"
    output_file = "Prezentare_Calitate_Aer.pptx"
    
    convert_md_to_pptx(md_file, output_file)
